"""PPTX -> ParsedDocument.

Uses python-pptx. Each slide's shapes are walked in document order (recursing into
group shapes, since PowerPoint lets authors group several shapes into one):
* title placeholder -> HeadingBlock (definitive; always wins over the heuristic below)
* other text frames -> paragraphs; bulleted/numbered paragraphs become list items
  (list_id per shape, list_level = paragraph indent level, ordered from a:buAutoNum
  vs a:buChar); inline formatting (bold/italic/underline/strike/super/subscript)
  is read from each run's a:rPr, falling back per-attribute to the paragraph's
  a:pPr/a:defRPr (OOXML's own resolution rule — see `_run_marks`), into `runs`
  on the block (mirrors docx's `Mark`/`InlineRun`; `text` stays the plain view)
* tables -> TableBlock with merges (a:gridSpan/rowSpan via python-pptx merge cells)
* pictures -> ImageBlock (locator = media part path; image_handler fetches bytes)
* embedded OLE objects (e.g. an embedded Excel sheet/Word doc) -> ImageBlock from
  the OOXML-mandated raster fallback preview (mc:Fallback/p:pic/blipFill/a:blip);
  skipped if a presentation has none (rare, but possible for some OLE types)
* speaker notes -> a ParagraphBlock per slide, spanned to
  ppt/notesSlides/notesSlideN.xml so consumers can tell notes from slide body

Slides that fake a title via a free-standing textbox (no TITLE placeholder at all —
common when a deck was built from scratch rather than a layout) are handled the same
way docx handles formatting-faked headings: candidate paragraphs are scored against
`PptxHeadingConfig` and, if a slide has no definitive title, its single best-scoring
candidate (if any clears the threshold) is promoted to a HeadingBlock. See
`heading_heuristics.py` for the shared pure-text clues (caps/title-case/date/caption)
and `docx_parser.HeadingConfig` for the sibling implementation this mirrors.

Locator: slides have no meaningful byte offset per shape, so blocks carry
Span(part="ppt/slides/slideN.xml", page=N) and leave byte offsets unset — the page
number is the reliable locator here (Karar B best-effort).
"""

from __future__ import annotations

import hashlib
from collections import Counter
from dataclasses import dataclass
from pathlib import Path

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from .base import (
    BaseParser, ParsedDocument, Span, Cell, text_cell,
    HeadingBlock, ParagraphBlock, TableBlock, ImageBlock, TableData, Merge,
    Mark, InlineRun, finalize_runs, runs_have_marks,
)
from .heading_heuristics import CAPTION, is_all_caps, is_title_case, looks_like_date, numbering_level

_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _clean_text(s: str) -> str:
    """python-pptx renders a soft line break (a:br) as \\x0b; flatten it to a space
    so it doesn't leak into the IR as a raw control character."""
    return s.replace("\x0b", " ")


def _bullet_kind(paragraph) -> str:
    """'ordered' | 'bullet' | 'none' | 'inherit' from the paragraph's a:pPr."""
    pPr = paragraph._p.find(qn("a:pPr"))
    if pPr is None:
        return "inherit"
    if pPr.find(qn("a:buNone")) is not None:
        return "none"
    if pPr.find(qn("a:buAutoNum")) is not None:
        return "ordered"
    if pPr.find(qn("a:buChar")) is not None:
        return "bullet"
    return "inherit"


def _table_data(table) -> TableData:
    nrows, ncols = len(table.rows), len(table.columns)
    cells: list[list[Cell | None]] = [[None] * ncols for _ in range(nrows)]
    merges: list[Merge] = []
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            if cell.is_spanned:
                continue  # covered by a merge origin -> stays None
            cells[r][c] = text_cell(_clean_text(cell.text))
            if cell.is_merge_origin and (cell.span_height > 1 or cell.span_width > 1):
                merges.append(Merge(row=r, col=c,
                                    rowspan=cell.span_height, colspan=cell.span_width))
    return TableData(n_rows=nrows, n_cols=ncols, cells=cells, merges=merges)


# ---------------------------------------------------------------------------
# Heading detection for slides that fake a title via a plain textbox.
#
# A real TITLE/CENTER_TITLE placeholder (PptxParser._is_title) is definitive and
# bypasses all of this. Only when a slide has NO such placeholder do we score the
# non-placeholder, non-list paragraphs on that slide and promote the single best
# one (if it clears the threshold) — a slide has one title, never several.
#
# Weights/threshold/gates are configurable via PptxHeadingConfig (never hardcoded).
# This is pptx's own knob set, calibrated separately from docx's HeadingConfig:
# the available signals differ (no w:outlineLvl/styles.xml baseline; bold/size/
# alignment come from a:rPr with a:pPr/a:defRPr fallback via python-pptx's
# `paragraph.font`; "isolation" only looks at neighboring paragraphs within the
# same shape, since paragraphs in different shapes aren't a text flow).
# ---------------------------------------------------------------------------


@dataclass
class PptxHeadingConfig:
    """Tunable knobs for scoring a free-standing textbox paragraph as a slide title.

    Nothing here is baked into the scoring logic — pass an instance to
    ``PptxParser(heading_config=...)`` to override any weight/threshold. Set a
    clue's weight to 0 to disable it; the ``enable_*`` flags turn whole gates off.
    """

    threshold: float = 6.0
    max_words: int = 14
    default_level: int = 1     # a slide has exactly one title -> always H1

    # positive clue weights (points added when the clue fires)
    bold: float = 3.0
    number: float = 2.0         # leading "1." / "I." / "Chapter" prefix
    caps: float = 2.0
    titlecase: float = 1.5      # Every Significant Word Capitalized (mixed case)
    titlecase_max_words: int = 7
    center: float = 1.5
    short: float = 2.0          # <= 7 words
    vshort: float = 1.0         # extra when <= 4 words (very short = short + vshort)
    isolation: float = 1.0      # blank paragraph both before and after, within the shape
    underline: float = 1.0
    no_punct: float = 1.0       # ends without terminal punctuation (near-always-on: low)
    large_font: float = 2.0     # font >= large_font_ratio x deck-wide baseline
    xlarge_font: float = 3.0    # font >= xlarge_font_ratio x deck-wide baseline

    # negative clue weights
    sentence: float = -3.0      # ends with . ? ! -> looks like a sentence
    italic: float = -2.0        # fully italic -> emphasis/citation, not a title

    # font-size clue, relative to the deck-wide baseline (most common explicit size)
    large_font_ratio: float = 1.15
    xlarge_font_ratio: float = 1.40

    # gates
    enable_font: bool = True
    enable_date_gate: bool = True


DEFAULT_PPTX_HEADING_CONFIG = PptxHeadingConfig()


def _run_marks(rPr, defRPr) -> tuple[Mark, ...]:
    """Semantic inline marks from a run's a:rPr, falling back per-attribute to
    the paragraph's a:pPr/a:defRPr — OOXML's own resolution rule for a run that
    doesn't set an attribute itself. Real corpus decks rely on this: e.g. a
    title's one run carries no b/sz at all, only the paragraph default does.
    """
    def attr(name: str) -> str | None:
        v = rPr.get(name) if rPr is not None else None
        return v if v is not None else (defRPr.get(name) if defRPr is not None else None)

    marks: list[Mark] = []
    if attr("b") in ("1", "true"):
        marks.append(Mark.BOLD)
    if attr("i") in ("1", "true"):
        marks.append(Mark.ITALIC)
    u = attr("u")
    if u and u != "none":
        marks.append(Mark.UNDERLINE)
    strike = attr("strike")
    if strike and strike != "noStrike":
        marks.append(Mark.STRIKE)
    baseline = attr("baseline")
    if baseline:
        try:
            v = int(baseline)
        except ValueError:
            v = 0
        if v > 0:
            marks.append(Mark.SUPERSCRIPT)
        elif v < 0:
            marks.append(Mark.SUBSCRIPT)
    return tuple(marks)


def _walk_pptx_para(para) -> list[InlineRun]:
    """Walk an a:p in document order -> InlineRuns (mirrors docx's `_walk_para`).

    a:br (soft line break) becomes an unmarked space, matching how a paragraph's
    flat text has always represented it (see the old `_clean_text` note this
    replaces for shape paragraphs). a:fld (a field, e.g. a slide-number
    placeholder) carries a cached a:t like a run and is read the same way.
    """
    pPr = para._p.find(qn("a:pPr"))
    defRPr = pPr.find(qn("a:defRPr")) if pPr is not None else None
    runs: list[InlineRun] = []
    for child in para._p:
        tag = child.tag
        if tag in (qn("a:r"), qn("a:fld")):
            rPr = child.find(qn("a:rPr"))
            t = child.find(qn("a:t"))
            runs.append(InlineRun(t.text or "" if t is not None else "",
                                  _run_marks(rPr, defRPr)))
        elif tag == qn("a:br"):
            runs.append(InlineRun(" "))
    return finalize_runs(runs)


def _para_font_size(para) -> float | None:
    """Largest resolved font size (pt) among a paragraph's text-bearing runs,
    falling back to a:pPr/a:defRPr per run (see `_run_marks`) — pptx decks set
    size either per-run or once on the paragraph default; mirrors docx's
    "largest explicit run size" heading clue.
    """
    pPr = para._p.find(qn("a:pPr"))
    defRPr = pPr.find(qn("a:defRPr")) if pPr is not None else None
    sizes: list[float] = []
    for run in para.runs:
        if not run.text.strip():
            continue
        rPr = run._r.find(qn("a:rPr"))
        sz = rPr.get("sz") if rPr is not None else None
        if sz is None and defRPr is not None:
            sz = defRPr.get("sz")
        if sz is not None:
            sizes.append(int(sz) / 100)   # a:rPr/@sz is in hundredths of a point
    if sizes:
        return max(sizes)
    if not para.runs and defRPr is not None and defRPr.get("sz") is not None:
        return int(defRPr.get("sz")) / 100
    return None


def _font_baseline(prs) -> float | None:
    """Deck-wide body baseline (pt): the most common resolved paragraph font size.

    Pptx has no single "Normal style" the way docx does; the most-common size
    across all text frames is a reasonable stand-in. None (font signal disabled)
    if no paragraph resolves to an explicit size (run rPr or paragraph defRPr).
    """
    sizes: Counter = Counter()

    def walk(shapes) -> None:
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        size = _para_font_size(para)
                        if size is not None:
                            sizes[round(size, 1)] += 1

    for slide in prs.slides:
        walk(slide.shapes)
    return sizes.most_common(1)[0][0] if sizes else None


def _pptx_heading_signals(runs: list[InlineRun], text: str, blank_before: bool,
                          blank_after: bool, baseline: float | None,
                          size: float | None, center: bool) -> dict:
    """Collect the transient heading clues for one non-placeholder paragraph."""
    textruns = [r for r in runs if r.text.strip()]
    bold = bool(textruns) and all(Mark.BOLD in r.marks for r in textruns)
    italic = bool(textruns) and all(Mark.ITALIC in r.marks for r in textruns)
    underline = bool(textruns) and all(Mark.UNDERLINE in r.marks for r in textruns)
    font_ratio = (size / baseline) if (size is not None and baseline) else 0.0
    stripped = text.rstrip()
    last = stripped[-1] if stripped else ""
    return {
        "text": text,
        "wc": len(text.split()),
        "caps": is_all_caps(text),
        "titlecase": is_title_case(text),
        "num": numbering_level(text),
        "bold": bold,
        "underline": underline,
        "italic": italic,
        "center": center,
        "font_ratio": font_ratio,
        "is_date": looks_like_date(text),
        "blank_before": blank_before,
        "blank_after": blank_after,
        "ends_sentence": last in ".?!",
        "no_terminal": bool(last) and last not in ".?!,;:",
        "sentences": sum(text.count(ch) for ch in ".?!"),
    }


def _classify_pptx_heading(sig: dict, cfg: PptxHeadingConfig) -> float | None:
    """Score one candidate; None if it's gated out (not heading-shaped at all)."""
    gate = (sig["wc"] > 0 and sig["sentences"] < 2
            and not CAPTION.match(sig["text"])
            and not (cfg.enable_date_gate and sig["is_date"])
            and (sig["num"] is not None or sig["wc"] <= cfg.max_words))
    if not gate:
        return None

    score = 0.0
    if sig["bold"]:
        score += cfg.bold
    if sig["num"] is not None:
        score += cfg.number
    if sig["caps"]:
        score += cfg.caps
    if sig["titlecase"] and sig["wc"] <= cfg.titlecase_max_words:
        score += cfg.titlecase
    if sig["center"]:
        score += cfg.center
    if sig["italic"]:
        score += cfg.italic
    if sig["wc"] <= 4:
        score += cfg.short + cfg.vshort
    elif sig["wc"] <= 7:
        score += cfg.short
    if sig["blank_before"] and sig["blank_after"]:
        score += cfg.isolation
    if sig["underline"]:
        score += cfg.underline
    if cfg.enable_font and sig["font_ratio"] >= cfg.xlarge_font_ratio:
        score += cfg.xlarge_font
    elif cfg.enable_font and sig["font_ratio"] >= cfg.large_font_ratio:
        score += cfg.large_font
    if sig["ends_sentence"]:
        score += cfg.sentence
    elif sig["no_terminal"]:
        score += cfg.no_punct
    return score


class PptxParser(BaseParser):
    extensions = (".pptx",)
    mimetypes = (_MIME,)
    fmt = "pptx"
    version = "pptx/0.1"

    def __init__(self, heading_config: PptxHeadingConfig | None = None) -> None:
        # Heading-detection knobs for formatting-faked slide titles; defaults are
        # a first calibration pass. Override with
        # PptxParser(PptxHeadingConfig(threshold=8, bold=4, ...)).
        self.heading_config = heading_config or DEFAULT_PPTX_HEADING_CONFIG

    def parse(self, raw_path: str | Path, doc_id: str) -> ParsedDocument:
        raw_path = Path(raw_path)
        data = raw_path.read_bytes()
        raw_sha256 = hashlib.sha256(data).hexdigest()

        prs = pptx.Presentation(raw_path)
        font_baseline = _font_baseline(prs)
        blocks: list = []
        bid = 0
        img_n = 0

        def next_id() -> str:
            nonlocal bid
            s = f"b{bid}"
            bid += 1
            return s

        for slide_no, slide in enumerate(prs.slides, start=1):
            span = Span(part=f"ppt/slides/slide{slide_no}.xml", page=slide_no)
            slide_start = len(blocks)
            candidates: list = []
            img_n = self._walk_shapes(slide.shapes, slide, slide_no, span,
                                      blocks, next_id, img_n, candidates)
            if candidates and not any(
                    isinstance(b, HeadingBlock) for b in blocks[slide_start:]):
                self._promote_pptx_heading(blocks, candidates,
                                           self.heading_config, font_baseline)
            if slide.has_notes_slide:
                notes = _clean_text(slide.notes_slide.notes_text_frame.text).strip()
                if notes:
                    notes_span = Span(part=f"ppt/notesSlides/notesSlide{slide_no}.xml",
                                      page=slide_no)
                    blocks.append(ParagraphBlock(id=next_id(), span=notes_span,
                                                 text=notes))

        return ParsedDocument(
            doc_id=doc_id,
            source_path=str(raw_path),
            fmt=self.fmt,
            raw_sha256=raw_sha256,
            mimetype=_MIME,
            page_count=len(prs.slides._sldIdLst),
            parser_version=self.version,
            blocks=blocks,
        )

    def _walk_shapes(self, shapes, slide, slide_no, span, blocks, next_id, img_n,
                     candidates) -> int:
        """Depth-first walk; recurses into group shapes so nested content isn't lost."""
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                img_n = self._walk_shapes(shape.shapes, slide, slide_no, span,
                                          blocks, next_id, img_n, candidates)
            elif shape.has_table:
                blocks.append(TableBlock(id=next_id(), span=span,
                                         table=_table_data(shape.table)))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_n += 1
                locator, mime = self._picture(shape, slide)
                blocks.append(ImageBlock(
                    id=next_id(), span=span, image_index=img_n,
                    locator=locator, mime=mime,
                    alt_text=shape.name or None))
            elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                locator, mime = self._picture(shape, slide)
                if locator:
                    img_n += 1
                    prog_id = None
                    try:
                        prog_id = shape.ole_format.prog_id
                    except (AttributeError, ValueError):
                        pass
                    alt = f"Embedded object ({prog_id})" if prog_id else (shape.name or None)
                    blocks.append(ImageBlock(
                        id=next_id(), span=span, image_index=img_n,
                        locator=locator, mime=mime, alt_text=alt))
            elif shape.has_text_frame:
                self._emit_text(shape, slide_no, span, blocks, next_id, candidates)
        return img_n

    _TITLE_PLACEHOLDERS = frozenset({"TITLE", "CENTER_TITLE", "VERTICAL_TITLE"})

    @classmethod
    def _is_title(cls, shape) -> bool:
        if not shape.is_placeholder:
            return False
        try:
            return shape.placeholder_format.type.name in cls._TITLE_PLACEHOLDERS
        except (AttributeError, ValueError):
            return False

    def _emit_text(self, shape, slide_no, span, blocks, next_id, candidates) -> None:
        is_title = self._is_title(shape)
        list_id = f"s{slide_no}_{shape.shape_id}"
        paras = list(shape.text_frame.paragraphs)
        run_lists = [_walk_pptx_para(p) for p in paras]
        texts = ["".join(r.text for r in rl) for rl in run_lists]
        for i, para in enumerate(paras):
            text = texts[i]
            if not text:
                continue
            runs = run_lists[i]
            marked = runs if runs_have_marks(runs) else []
            if is_title:
                blocks.append(HeadingBlock(id=next_id(), span=span,
                                           text=text, level=1, runs=marked))
                continue
            kind = _bullet_kind(para)
            is_list = (kind in ("ordered", "bullet")
                       or (kind != "none" and (shape.is_placeholder or para.level > 0)))
            if is_list:
                ordered = True if kind == "ordered" else (False if kind == "bullet" else None)
                blocks.append(ParagraphBlock(
                    id=next_id(), span=span, text=text, runs=marked,
                    list_id=list_id, list_level=para.level, list_ordered=ordered))
            else:
                pb = ParagraphBlock(id=next_id(), span=span, text=text, runs=marked)
                blocks.append(pb)
                if not shape.is_placeholder:
                    blank_before = i > 0 and not texts[i - 1]
                    blank_after = i < len(texts) - 1 and not texts[i + 1]
                    candidates.append((pb, para, runs, blank_before, blank_after))

    @staticmethod
    def _promote_pptx_heading(blocks: list, candidates: list,
                              cfg: PptxHeadingConfig, baseline: float | None) -> None:
        """Promote the single best-scoring candidate paragraph on a slide to a
        HeadingBlock in place, if any clears the threshold. A slide has one title."""
        best_pb, best_score = None, cfg.threshold - 1e-9
        for pb, para, runs, blank_before, blank_after in candidates:
            size = _para_font_size(para)
            center = para.alignment == PP_ALIGN.CENTER
            sig = _pptx_heading_signals(runs, pb.text, blank_before, blank_after,
                                        baseline, size, center)
            score = _classify_pptx_heading(sig, cfg)
            if score is not None and score > best_score:
                best_pb, best_score = pb, score
        if best_pb is None:
            return
        for i, b in enumerate(blocks):
            if b is best_pb:
                hruns = b.runs if runs_have_marks(b.runs) else []
                blocks[i] = HeadingBlock(id=best_pb.id, span=best_pb.span,
                                         text=best_pb.text, level=cfg.default_level,
                                         runs=hruns)
                break

    @staticmethod
    def _picture(shape, slide) -> tuple[dict, str | None]:
        locator: dict = {}
        mime = None
        try:
            mime = shape.image.content_type
        except (AttributeError, ValueError):
            pass
        blip = shape._element.find(".//" + qn("a:blip"))
        if blip is not None:
            embed = blip.get(qn("r:embed"))
            try:
                part = slide.part.related_part(embed)
                locator = {"part": str(part.partname).lstrip("/")}
                if mime is None:
                    mime = part.content_type
            except KeyError:
                pass
        return locator, mime
