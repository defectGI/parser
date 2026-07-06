"""PPTX extraction fixes:

1. Group shapes (p:grpSp) were not walked -> nested text/pictures/tables were
   silently dropped.
2. Embedded OLE objects (e.g. an embedded Excel sheet) were not handled at all
   -> the OOXML-mandated raster fallback preview (p:pic/blipFill/a:blip) is now
   captured as an ImageBlock.
3. Speaker notes were never read -> now emitted as a ParagraphBlock spanned to
   ppt/notesSlides/notesSlideN.xml.
4. `_is_title` matched "TITLE" as a substring, so SUBTITLE placeholders (whose
   enum name contains "TITLE") were wrongly promoted to HeadingBlocks.
5. Slides that fake a title via a free-standing textbox (no TITLE placeholder
   at all) never produced a HeadingBlock -> scored heuristic promotion, mirroring
   docx's formatting-faked-heading detection (see PptxHeadingConfig).
6. Inline formatting (bold/italic/underline/strike/super/subscript) was never
   read -> `runs` was always empty for pptx blocks. Now read per-run from a:rPr,
   falling back to the paragraph's a:pPr/a:defRPr for whichever attribute a run
   doesn't set itself (both patterns occur in real decks).
"""

from __future__ import annotations

import io
import struct
import zlib

from pptx import Presentation
from pptx.enum.shapes import PROG_ID
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from parsers.base import HeadingBlock, ImageBlock, Mark, ParagraphBlock, ParsedDocument
from parsers.pptx_parser import PptxParser


def _png_1x1() -> bytes:
    def chunk(typ: bytes, data: bytes) -> bytes:
        body = typ + data
        return (struct.pack(">I", len(data)) + body
                + struct.pack(">I", zlib.crc32(body) & 0xFFFFFFFF))
    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    idat = zlib.compress(b"\x00\xff\x00\x00")
    return (b"\x89PNG\r\n\x1a\n"
            + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b""))


def test_text_inside_group_shape_is_captured(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb.text_frame.text = "Grouped Label"
    slide.shapes.add_group_shape([tb])
    p = tmp_path / "grouped.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "grouped")
    paras = [b for b in doc.blocks if isinstance(b, ParagraphBlock)]
    assert [b.text for b in paras] == ["Grouped Label"]


def test_speaker_notes_become_a_paragraph_block(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "Talking points for slide one."
    p = tmp_path / "notes.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "notes")
    notes = [b for b in doc.blocks
             if isinstance(b, ParagraphBlock) and "notesSlide" in (b.span.part or "")]
    assert len(notes) == 1
    assert notes[0].text == "Talking points for slide one."
    assert notes[0].span.part == "ppt/notesSlides/notesSlide1.xml"


def test_slide_without_notes_emits_no_notes_block(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    p = tmp_path / "nonotes.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "nonotes")
    assert not any("notesSlide" in (b.span.part or "") for b in doc.blocks)


def test_ole_object_fallback_preview_becomes_an_image(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_ole_object(
        io.BytesIO(b"fake xlsx bytes" * 10), PROG_ID.XLSX,
        Inches(1), Inches(3),
        icon_file=io.BytesIO(_png_1x1()), icon_width=Inches(1), icon_height=Inches(1))
    p = tmp_path / "ole.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "ole")
    imgs = doc.images()
    assert len(imgs) == 1
    assert imgs[0].locator["part"].startswith("ppt/media/")
    assert "Excel.Sheet.12" in imgs[0].alt_text


def _bold_run(shape, text, *, size=None, bold=None, center=False):
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if center:
        p.alignment = PP_ALIGN.CENTER
    return p


def test_subtitle_placeholder_is_not_promoted_to_heading(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide: title + subtitle
    slide.placeholders[0].text_frame.text = "Chapter 4"
    slide.placeholders[1].text_frame.text = "Selected Questions"
    p = tmp_path / "subtitle.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "subtitle")
    headings = [b for b in doc.blocks if isinstance(b, HeadingBlock)]
    assert [h.text for h in headings] == ["Chapter 4"]


def test_freeform_title_textbox_is_promoted_to_heading(tmp_path):
    """No TITLE placeholder at all (deck built from plain textboxes, like
    BlueYonder_Presentation.pptx in the real corpus) -> the bold/centered/short
    textbox should still be recognized as the slide's title."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(6), Inches(1))
    _bold_run(title_box, "Quarterly Sales Review", size=32, bold=True, center=True)

    body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    _bold_run(body_box, "This section walks through the quarterly numbers in detail.")

    p = tmp_path / "freeform.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "freeform")
    headings = [b for b in doc.blocks if isinstance(b, HeadingBlock)]
    paras = [b for b in doc.blocks if isinstance(b, ParagraphBlock)]
    assert [h.text for h in headings] == ["Quarterly Sales Review"]
    assert [pb.text for pb in paras] == [
        "This section walks through the quarterly numbers in detail."]


def test_caption_like_bold_line_is_not_promoted(tmp_path):
    """Bold + centered + short would otherwise score above threshold, but a
    figure/table caption is gated out explicitly (never a heading)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    _bold_run(box, "Table 1: Revenue by Quarter", size=24, bold=True, center=True)
    p = tmp_path / "caption.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "caption")
    assert not any(isinstance(b, HeadingBlock) for b in doc.blocks)


def test_only_the_best_candidate_is_promoted_per_slide(tmp_path):
    """Two textboxes both clear the threshold; only the single best-scoring one
    (a slide has exactly one title) becomes a HeadingBlock."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    strong = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(6), Inches(1))
    _bold_run(strong, "Executive Summary", size=32, bold=True, center=True)

    weak = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    _bold_run(weak, "Overview", size=20, bold=True, center=True)

    p = tmp_path / "twocandidates.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "twocandidates")
    headings = [b for b in doc.blocks if isinstance(b, HeadingBlock)]
    assert [h.text for h in headings] == ["Executive Summary"]
    paras = [b for b in doc.blocks if isinstance(b, ParagraphBlock)]
    assert [pb.text for pb in paras] == ["Overview"]


def test_mid_paragraph_run_marks_are_captured(tmp_path):
    """A paragraph mixing a plain run and a bold+italic run keeps one plain
    `text` view but records the per-run marks in `runs`."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    para = box.text_frame.paragraphs[0]
    r1 = para.add_run()
    r1.text = "Please note: "
    r2 = para.add_run()
    r2.text = "handle with care"
    r2.font.bold = True
    r2.font.italic = True
    p = tmp_path / "runs.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "runs")
    pb = next(b for b in doc.blocks if isinstance(b, ParagraphBlock))
    assert pb.text == "Please note: handle with care"
    assert [(r.text, r.marks) for r in pb.runs] == [
        ("Please note: ", ()),
        ("handle with care", (Mark.BOLD, Mark.ITALIC)),
    ]


def test_paragraph_defrpr_fallback_marks_an_unmarked_run(tmp_path):
    """A run with no rPr of its own inherits bold from the paragraph's
    a:pPr/a:defRPr (python-pptx's `paragraph.font`) — the same fallback the
    title heuristic already relies on. Text is a full sentence so the heading
    heuristic doesn't also promote it (this test is only about `runs`)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    para = box.text_frame.paragraphs[0]
    para.font.bold = True
    run = para.add_run()
    run.text = "This entire sentence inherits bold from the paragraph default."
    p = tmp_path / "defrpr.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "defrpr")
    b = next(b for b in doc.blocks if b.text.startswith("This entire sentence"))
    assert [(r.text, r.marks) for r in b.runs] == [(b.text, (Mark.BOLD,))]


def test_soft_line_break_becomes_an_unmarked_space(tmp_path):
    """a:br (Shift+Enter inside one paragraph) used to leak into the IR as a
    raw \\x0b control character; it's now a plain space in both `text` and `runs`."""
    from pptx.oxml.ns import qn as _qn
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    para = box.text_frame.paragraphs[0]
    r1 = para.add_run()
    r1.text = "Line one"
    para._p.append(para._p.makeelement(_qn("a:br"), {}))
    r2 = para.add_run()
    r2.text = "Line two"
    p = tmp_path / "linebreak.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "linebreak")
    pb = next(b for b in doc.blocks if isinstance(b, ParagraphBlock))
    assert pb.text == "Line one Line two"
    assert "\x0b" not in pb.text


def test_pptx_roundtrip_with_runs(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "bold text"
    run.font.bold = True
    p = tmp_path / "rt.pptx"
    prs.save(str(p))

    doc = PptxParser().parse(p, "rt")
    restored = ParsedDocument.from_json(doc.to_json())
    assert restored.to_dict() == doc.to_dict()
